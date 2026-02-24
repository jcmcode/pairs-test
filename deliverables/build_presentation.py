"""
Build presentation.pptx for Queen's Quant Club judges.
Run: python deliverables/build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colours
QUEENS_BLUE = RGBColor(0, 36, 82)
QUEENS_RED = RGBColor(185, 17, 55)
DARK_GRAY = RGBColor(60, 60, 60)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(39, 174, 96)
RED = RGBColor(231, 76, 60)
BLACK = RGBColor(0, 0, 0)
ACCENT_BLUE = RGBColor(230, 240, 250)
ACCENT_GREEN = RGBColor(230, 250, 235)
ACCENT_RED = RGBColor(255, 235, 235)
GOLD = RGBColor(250, 189, 0)
SUBTLE_GRAY = RGBColor(100, 100, 100)
HEADER_BLUE = RGBColor(180, 200, 220)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_counter = [0]  # mutable counter for slide numbers


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number(slide):
    slide_counter[0] += 1
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_counter[0])
    p.font.size = Pt(10)
    p.font.color.rgb = SUBTLE_GRAY
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.RIGHT


def add_speaker_notes(slide, notes_text):
    slide.notes_slide.notes_text_frame.text = notes_text


def text_box(slide, left, top, width, height, text, font_size=18,
             bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def bullet_slide(slide, left, top, width, height, items, font_size=16,
                 color=DARK_GRAY, bold_items=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
        if bold_items and i in bold_items:
            p.font.bold = True
    return tf


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = 'Calibri'
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = QUEENS_BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = 'Calibri'
                paragraph.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table


def section_header(slide, title, subtitle=""):
    add_bg(slide, QUEENS_BLUE)
    text_box(slide, 1.5, 2.5, 10, 1.5, title, font_size=40, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT)
    if subtitle:
        text_box(slide, 1.5, 4.2, 10, 1, subtitle, font_size=20, color=HEADER_BLUE,
                 alignment=PP_ALIGN.LEFT)
    add_slide_number(slide)


def key_idea_box(slide, left, top, width, height, text, border_color=QUEENS_BLUE, bg_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = border_color
    p.font.name = 'Calibri'
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return tf


def pipeline_box(slide, left, top, width, height, label, sublabel="",
                 bg_color=QUEENS_BLUE, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    if sublabel:
        p2 = tf.add_paragraph()
        p2.text = sublabel
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(200, 220, 240)
        p2.font.name = 'Calibri'
        p2.alignment = PP_ALIGN.CENTER


def pipeline_arrow(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(left), Inches(top), Inches(0.4), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = QUEENS_RED
    shape.line.fill.background()


# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, QUEENS_BLUE)

text_box(slide, 1.5, 1.5, 10, 1.5,
         "Clustering-Driven Pair Discovery\nFrom Semiconductors to Cross-Sector Markets",
         font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

text_box(slide, 1.5, 3.8, 10, 0.8,
         "Unsupervised Learning for Transient Correlation Detection Across 142 Tickers",
         font_size=20, color=HEADER_BLUE, alignment=PP_ALIGN.LEFT)

text_box(slide, 1.5, 5.5, 10, 0.6,
         "Queen's Quant Club  |  February 2026",
         font_size=18, color=GOLD, alignment=PP_ALIGN.LEFT)

add_slide_number(slide)
add_speaker_notes(slide, "Title slide. This project explores using unsupervised clustering to discover pairs for trading, moving beyond traditional cointegration-based approaches.")


# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 5, 0.8, "Agenda", font_size=32, bold=True, color=QUEENS_BLUE)

items = [
    "1.  Motivation: Why clustering for pair discovery?",
    "2.  Methodology: Features, algorithms, pipeline",
    "3.  Phase 1: Semiconductor proof of concept",
    "4.  Phase 2: Cross-sector expansion (142 tickers, 5 sectors)",
    "5.  5-Test Validation Framework",
    "6.  Enhanced Backtesting & Walk-Forward",
    "7.  Statistical Significance",
    "8.  Key Insights & Future Work",
]
bullet_slide(slide, 1.2, 1.5, 10, 5, items, font_size=20, color=DARK_GRAY)
add_slide_number(slide)
add_speaker_notes(slide, "Overview of the presentation structure. 8 sections covering the full research arc from motivation to results.")


# ============================================================================
# SLIDE 3: The Problem
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 8, 0.8, "The Problem: Finding Pairs to Trade", font_size=28, bold=True, color=QUEENS_BLUE)

# Left column - Traditional
text_box(slide, 0.8, 1.5, 5.5, 0.6, "Traditional Approach", font_size=20, bold=True, color=QUEENS_RED)
bullet_slide(slide, 0.8, 2.2, 5.5, 3, [
    "Screen all C(N,2) combinations for cointegration",
    "For 142 stocks: 10,011 pairs to test",
    "High false discovery rate from multiple testing",
    "No insight into WHY pairs are related",
], font_size=16)

# Right column - Our approach
text_box(slide, 7, 1.5, 5.5, 0.6, "Our Approach", font_size=20, bold=True, color=GREEN)
bullet_slide(slide, 7, 2.2, 5.5, 3, [
    "Let clustering discover similar-behaving stocks",
    "Validate with a 5-test scored framework",
    "Compare across 3 different algorithms",
    "142 tickers, 5 sectors, 12 months hourly data",
], font_size=16)

# Key idea box
key_idea_box(slide, 2.5, 5.3, 8, 1.3,
             "Key Idea: Stocks that consistently cluster together based on\n"
             "real-time feature similarity are candidates for pairs trading.")

add_slide_number(slide)
add_speaker_notes(slide, "Traditional pairs trading screens all N-choose-2 pairs for cointegration, which is brute-force and gives no economic intuition. Our approach uses clustering to discover which stocks behave similarly, then validates those relationships statistically.")


# ============================================================================
# SLIDE 4: Pipeline Overview (visual diagram)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Pipeline Overview", font_size=28, bold=True, color=QUEENS_BLUE)

# Pipeline boxes with arrows
stages = [
    ("Data", "Yahoo Finance\nHourly OHLCV", 0.5),
    ("Features", "9 per ticker\nper timestamp", 2.4),
    ("Normalize", "StandardScaler\n+ PCA (90%)", 4.3),
    ("Cluster", "OPTICS\nper timestamp", 6.2),
    ("Discover", "Co-clustering\npairs", 8.1),
    ("Validate", "5-test scored\nframework", 10.0),
]

arrow_positions = [1.95, 3.85, 5.75, 7.65, 9.55]

y_pipe = 1.6
for label, sublabel, x in stages:
    pipeline_box(slide, x, y_pipe, 1.5, 0.9, label, sublabel)

for x in arrow_positions:
    pipeline_arrow(slide, x, y_pipe + 0.3)

# Backtest box at end
pipeline_box(slide, 11.3, y_pipe, 1.3, 0.9, "Backtest", "Walk-forward\n+ Kalman", bg_color=QUEENS_RED)
pipeline_arrow(slide, 11.0, y_pipe + 0.3)

# Features table
text_box(slide, 0.8, 3.0, 10, 0.5, "Feature Set (9 features per ticker per timestamp)",
         font_size=18, bold=True, color=QUEENS_BLUE)

headers = ["Category", "Feature", "Window"]
rows = [
    ["Volatility", "Short-term rolling \u03c3", "50h (~1 week)"],
    ["Market exposure", "Rolling \u03b2 to S&P 500", "50h"],
    ["Sector exposure", "Rolling \u03b2 to leave-one-out sector avg", "50h"],
    ["Momentum", "RSI (70h) + 5-hour return", "70h / Instant"],
    ["Regime shift", "\u0394\u03c3, \u0394\u03b2_SPX, \u0394\u03b2_sector (short \u2212 medium)", "50h vs 147h"],
    ["Returns", "Hourly return", "Instantaneous"],
]
add_table(slide, 1.5, 3.6, 10, 3.2, headers, rows, col_widths=[2.5, 5.0, 2.5])

add_slide_number(slide)
add_speaker_notes(slide, "The pipeline runs per-timestamp: compute 9 features for each ticker, standardize and reduce with PCA, cluster with OPTICS, then identify pairs from co-clustering. Features capture volatility, market/sector beta, momentum, and regime shifts. The regime shift features (short minus medium window) are the key signal for detecting WHEN relationships form.")


# ============================================================================
# SLIDE 5: Three Algorithms
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Three Algorithms, One Question",
         font_size=28, bold=True, color=QUEENS_BLUE)

text_box(slide, 0.8, 1.2, 11, 0.6,
         "If different algorithms with different inductive biases find the same pairs, those relationships are robust.",
         font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

headers = ["", "OPTICS", "DBSCAN", "KMeans"]
rows = [
    ["Type", "Density-based", "Density-based", "Centroid-based"],
    ["k required?", "No (auto)", "No (auto)", "Yes (silhouette)"],
    ["Noise handling", "Native (-1 label)", "Native", "Post-hoc filter"],
    ["Avg noise rate", "58%", "28%", "~10%"],
    ["Avg clusters/ts", "3.3", "1.5", "3.1"],
    ["OOS stability (r)", "0.671", "0.832", "0.627"],
]
add_table(slide, 1.5, 2.0, 10, 3.5, headers, rows, col_widths=[2.8, 2.4, 2.4, 2.4])

# Noise-adjusted frequency formula
key_idea_box(slide, 1.0, 5.8, 11, 1.2,
             "Noise-Adjusted Frequency:  f(A,B) = co-cluster count / timestamps where BOTH are non-noise\n"
             "With 58% noise rate, raw frequency drastically understates true relationship strength.")

add_slide_number(slide)
add_speaker_notes(slide, "OPTICS selected as primary algorithm for Phase 2: best at detecting variable-density clusters and transient formations. The 58% noise rate is by design -- OPTICS only forms clusters when there is genuine feature similarity. Noise-adjusted frequency corrects the denominator to count only timestamps where both tickers are in real clusters.")


# ============================================================================
# SLIDE 6: Phase 1 - Consensus Pairs
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Phase 1: Semiconductors", "40 tickers  |  3 algorithms  |  8 consensus pairs")
add_speaker_notes(slide, "Phase 1 was the proof of concept on 40 hand-picked semiconductor tickers. All three algorithms were run and compared.")


# ============================================================================
# SLIDE 7: Consensus Pairs + Subsector Discovery
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "8 Consensus Pairs Across All Algorithms",
         font_size=28, bold=True, color=QUEENS_BLUE)

text_box(slide, 0.8, 1.1, 11, 0.4,
         "Pairs appearing in every algorithm's top-20 by co-clustering frequency:",
         font_size=14, color=SUBTLE_GRAY)

headers = ["Pair", "OPTICS", "KMeans", "DBSCAN", "Subsector"]
rows = [
    ["KLAC\u2013LRCX", "0.378", "0.785", "0.851", "Equipment"],
    ["QRVO\u2013SWKS", "0.372", "0.888", "0.828", "RF / 5G"],
    ["AMAT\u2013LRCX", "0.363", "0.824", "0.845", "Equipment"],
    ["AMAT\u2013KLAC", "0.340", "0.807", "0.814", "Equipment"],
    ["ADI\u2013NXPI", "0.279", "0.816", "0.900", "Analog"],
    ["ADI\u2013TXN", "0.278", "0.887", "0.817", "Analog"],
    ["NXPI\u2013TXN", "0.223", "0.819", "0.806", "Analog"],
    ["ADI\u2013SWKS", "0.216", "0.830", "0.830", "Analog/RF"],
]
add_table(slide, 0.8, 1.5, 5.8, 4.0, headers, rows, col_widths=[1.4, 1.0, 1.0, 1.0, 1.2])

# Subsector discovery (right side)
text_box(slide, 7.0, 1.1, 5.5, 0.4, "Discovered Industry Structure",
         font_size=18, bold=True, color=QUEENS_BLUE)

text_box(slide, 7.0, 1.5, 5.5, 0.5,
         "No fundamental data as input \u2014 only price-derived features:",
         font_size=13, color=SUBTLE_GRAY)

clusters = [
    ("Equipment", "AMAT, LRCX, KLAC"),
    ("RF / 5G", "QRVO, SWKS, QCOM"),
    ("Analog", "ADI, NXPI, TXN, MCHP"),
    ("EDA", "CDNS, SNPS"),
]

for i, (name, tickers) in enumerate(clusters):
    y = 2.1 + i * 0.8
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(7.0), Inches(y), Inches(5.5), Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = QUEENS_BLUE
    text_box(slide, 7.2, y + 0.08, 1.8, 0.4, name, font_size=15, bold=True, color=QUEENS_BLUE)
    text_box(slide, 9.2, y + 0.08, 3, 0.4, tickers, font_size=14, color=DARK_GRAY)

# Bottom line
text_box(slide, 0.8, 5.8, 11.5, 0.6,
         "14/26 tradeable pairs profitable OOS (54%). Top Sharpe: NXPI\u2013ON 3.55, CDNS\u2013SNPS 2.08.",
         font_size=14, color=DARK_GRAY)

key_idea_box(slide, 0.8, 6.3, 11.5, 0.8,
             "Data-driven clustering converges with known industry structure \u2192 discovered relationships are real, not noise.")

add_slide_number(slide)
add_speaker_notes(slide, "The 8 consensus pairs all map to known semiconductor subsectors. Equipment (AMAT/LRCX/KLAC share fab customers), RF/5G (smartphone demand), Analog (industrial/auto), EDA (chip design duopoly). 54% of tradeable pairs were profitable OOS, but trade counts are low (2-5 per pair).")


# ============================================================================
# SLIDE 8: Phase 2 Section Header
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Phase 2: Cross-Sector", "142 tickers  |  5 sectors  |  3,643 pairs scored")
add_speaker_notes(slide, "Phase 2 extends the framework to 142 systematically screened tickers across Technology, Healthcare, Energy, Financial Services, and Industrials.")


# ============================================================================
# SLIDE 9: Universe Construction
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Universe Construction: 3-Layer Pipeline",
         font_size=28, bold=True, color=QUEENS_BLUE)

# Three layers
layers = [
    ("Layer 1: Liquidity", "US-listed, market cap > $2B, volume > 5M/day, price > $5"),
    ("Layer 2: Sector", "Technology, Healthcare, Energy, Financial Services, Industrials"),
    ("Layer 3: Quality", "Positive EBITDA + data quality (>80% coverage, <5% zero prices)"),
]
for i, (name, desc) in enumerate(layers):
    y = 1.4 + i * 1.0
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1.2), Inches(y), Inches(10.8), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = QUEENS_BLUE
    text_box(slide, 1.5, y + 0.1, 3.5, 0.5, name, font_size=17, bold=True, color=QUEENS_BLUE)
    text_box(slide, 5.2, y + 0.1, 6.5, 0.5, desc, font_size=15, color=DARK_GRAY)

# Key design decisions
text_box(slide, 0.8, 4.7, 11, 0.5, "Key Design Decisions", font_size=20, bold=True, color=QUEENS_BLUE)
bullet_slide(slide, 1.2, 5.3, 10, 2, [
    "All 142 tickers clustered together (not per-sector) \u2014 cross-sector pairs discovered organically",
    "7,836 co-clustering pairs identified; 3,643 scored by 5-test framework at 0.08 noise-adjusted frequency threshold",
    "Frequency threshold lowered from 0.15 (Phase 1) to 0.08 \u2014 at 0.15, zero cross-sector pairs survive",
], font_size=15, color=DARK_GRAY)

add_slide_number(slide)
add_speaker_notes(slide, "The three-layer pipeline systematically screens stocks. All tickers are clustered together in one combined universe -- cross-sector pairs emerge naturally. 7,836 total co-clustering pairs discovered; 3,643 pass the 0.08 noise-adjusted frequency threshold and are evaluated by the 5-test validation framework.")


# ============================================================================
# SLIDE 10: Why Cointegration Fails + 5-Test Framework
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Why Cointegration Fails & The 5-Test Solution",
         font_size=28, bold=True, color=QUEENS_BLUE)

# Problem box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3),
                                Inches(5.5), Inches(2.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.color.rgb = QUEENS_RED
text_box(slide, 1.1, 1.4, 5, 0.4, "The Problem", font_size=18, bold=True, color=QUEENS_RED)
bullet_slide(slide, 1.1, 1.9, 5, 1.5, [
    "Engle-Granger cointegration: 0% pass rate",
    "Tests for long-run equilibrium \u2014 wrong model",
    "Yet Hurst passes 90.5%, Half-life 82.3%",
    "Spreads DO mean-revert, just not permanently",
], font_size=14, color=DARK_GRAY)

# Solution box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.3),
                                Inches(5.5), Inches(2.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_GREEN
shape.line.color.rgb = GREEN
text_box(slide, 7.3, 1.4, 5, 0.4, "The Solution: 5-Test Scored Framework", font_size=18, bold=True, color=GREEN)
bullet_slide(slide, 7.3, 1.9, 5, 1.5, [
    "Each test scores 0 or 1 (total 0\u20135)",
    "Strong (4\u20135), Moderate (3), Weak (2), Fail (<2)",
    "59.0% of pairs are tradeable (score \u2265 3)",
    "vs 0% with cointegration!",
], font_size=14, color=DARK_GRAY)

# Framework table
headers = ["Test", "Threshold", "Pass Rate", "Role"]
rows = [
    ["ADF on spread", "p < 0.10", "38.3%", "Stationarity (hard)"],
    ["Half-life", "5\u201360 days", "82.3%", "Mean-reversion speed"],
    ["Hurst exponent", "H < 0.5", "90.5%", "Mean-reverting (easy)"],
    ["Variance ratio", "Reject RW @10%", "21.8%", "Short-run MR (hardest)"],
    ["Rolling corr", "Stability > 0.5", "34.7%", "Return corr persistence"],
]
add_table(slide, 1.5, 3.6, 10, 3.2, headers, rows, col_widths=[2.2, 2.0, 1.5, 4.3])

add_slide_number(slide)
add_speaker_notes(slide, "Cointegration assumes a permanent equilibrium relationship (Engle-Granger 1987). Transient correlations violate this assumption -- they form and dissolve over weeks. But the underlying spreads DO mean-revert (90.5% Hurst < 0.5). The 5-test framework replaces binary cointegration with a graduated score. ADF and variance ratio are the hard differentiators; Hurst and half-life are easy passes.")


# ============================================================================
# SLIDE 11: Score Distribution
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Score Distribution (3,643 Pairs)",
         font_size=28, bold=True, color=QUEENS_BLUE)

# Score distribution table (left)
headers = ["Classification", "Count", "%"]
rows = [
    ["Strong (4\u20135)", "605", "16.6%"],
    ["Moderate (3)", "1,543", "42.4%"],
    ["Weak (2)", "1,200", "32.9%"],
    ["Fail (0\u20131)", "295", "8.1%"],
    ["Tradeable (\u2265 3)", "2,148", "59.0%"],
]
add_table(slide, 0.8, 1.3, 5.2, 3.2, headers, rows, col_widths=[2.2, 1.5, 1.5])

# Intra vs Cross comparison (right)
text_box(slide, 7, 1.0, 5.5, 0.5, "Intra-Sector vs Cross-Sector", font_size=18, bold=True, color=QUEENS_BLUE)
headers = ["Type", "Pairs", "Avg Score", "Tradeable %", "Avg Sharpe"]
rows = [
    ["Intra-sector", "1,442", "2.77", "64.1%", "1.95"],
    ["Cross-sector", "2,201", "2.62", "55.6%", "1.59"],
]
add_table(slide, 6.8, 1.6, 5.8, 1.5, headers, rows, col_widths=[1.5, 0.9, 1.1, 1.2, 1.1])

# Sector matrix
text_box(slide, 7, 3.5, 5.5, 0.5, "Top Sector Combinations", font_size=18, bold=True, color=QUEENS_BLUE)
headers = ["Sectors", "Pairs", "Avg Score"]
rows = [
    ["Technology\u2013Technology", "880", "2.64"],
    ["Energy\u2013Energy", "270", "3.09"],
    ["Industrials\u2013Technology", "243", "2.79"],
    ["Healthcare\u2013Technology", "197", "2.57"],
    ["Healthcare\u2013Healthcare", "183", "2.88"],
]
add_table(slide, 6.8, 4.1, 5.8, 2.8, headers, rows, col_widths=[2.8, 1.2, 1.8])

key_idea_box(slide, 0.8, 5.0, 5.2, 1.5,
             "Intra-sector pairs have higher scores on average,\n"
             "but cross-sector pairs produce the best individual\n"
             "Sharpe ratio (2.56 for ADT\u2013NOW).")

add_slide_number(slide)
add_speaker_notes(slide, "59.0% of all scored pairs are tradeable (score >= 3). Energy-Energy pairs have the highest average score (3.09), reflecting tight economic coupling in the oil & gas sector. Intra-sector pairs outperform cross-sector on average, but cross-sector pairs can produce strong individual results.")


# ============================================================================
# SLIDE 12: Cross-Sector Discoveries
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Cross-Sector Discoveries",
         font_size=28, bold=True, color=QUEENS_BLUE)

# RTX-SHEL
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3),
                                Inches(5.5), Inches(2.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.color.rgb = QUEENS_BLUE
text_box(slide, 1.1, 1.4, 5, 0.4, "RTX\u2013SHEL (Score 4, Strong)", font_size=18, bold=True, color=QUEENS_BLUE)
bullet_slide(slide, 1.1, 1.9, 5, 1.2, [
    "Industrials / Energy cross-sector link",
    "Defense contractor \u2194 Energy major",
    "Enhanced PnL: +10.07, Kalman: +13.87",
], font_size=14, color=DARK_GRAY)

# APLD Cluster
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.3),
                                Inches(5.5), Inches(2.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.color.rgb = QUEENS_BLUE
text_box(slide, 7.3, 1.4, 5, 0.4, "Crypto Mining Cluster", font_size=18, bold=True, color=QUEENS_BLUE)
bullet_slide(slide, 7.3, 1.9, 5, 1.2, [
    "APLD / HUT / CIFR / IREN across Tech & FinServ",
    "APLD-HUT: Score 4 (Strong)",
    "Meaningful economic link: shared infrastructure",
], font_size=14, color=DARK_GRAY)

# Top pairs table
text_box(slide, 0.8, 3.6, 10, 0.5, "Top Pairs (by score then Sharpe, min 5 OOS trades)", font_size=18, bold=True, color=QUEENS_BLUE)
headers = ["Type", "Pair", "Score", "ADF p", "Hurst", "Trades", "Sharpe"]
rows = [
    ["Cross", "ADT\u2013NOW", "4", "0.0002", "0.11", "5", "2.56"],
    ["Intra", "APH\u2013TSM", "4", "0.529", "0.49", "5", "1.93"],
    ["Cross", "CRBG\u2013NOW", "4", "0.323", "0.33", "5", "1.64"],
    ["Intra", "CRBG\u2013WU", "3", "0.178", "0.23", "5", "1.96"],
    ["Cross", "DXCM\u2013NOK", "3", "0.312", "0.29", "5", "0.57"],
]
add_table(slide, 1.5, 4.2, 10, 2.8, headers, rows, col_widths=[1.0, 2.0, 1.0, 1.2, 1.2, 1.2, 1.2])

add_slide_number(slide)
add_speaker_notes(slide, "RTX-SHEL is the most surprising cross-sector discovery: a defense contractor and energy major linked through geopolitical risk exposure. The crypto mining cluster (APLD/HUT/CIFR/IREN) spans Technology and Financial Services but shares common infrastructure and Bitcoin price exposure. ADT-NOW is the top-performing cross-sector pair with Sharpe 2.56 on 5 trades.")


# ============================================================================
# SLIDE 13: Enhanced Backtesting Section
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Enhanced Backtesting", "Adaptive z-scores  |  Kalman hedge  |  Walk-forward validation")
add_speaker_notes(slide, "Phase 2 enhanced the backtesting methodology with three key improvements over Phase 1's simple static z-score approach.")


# ============================================================================
# SLIDE 14: Three-Strategy Comparison
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Three-Strategy Comparison",
         font_size=28, bold=True, color=QUEENS_BLUE)

text_box(slide, 0.8, 1.1, 11, 0.5,
         "50 tradeable pairs backtested (of 2,148 scoring \u2265 3; optimizer requires \u2265 5 calibration trades).",
         font_size=13, color=SUBTLE_GRAY)

headers = ["Strategy", "Description", "Trades/Pair", "Profitable", "Avg PnL"]
rows = [
    ["Baseline", "Static z=2.0, OLS, no costs", "1.9", "54%", "0.123"],
    ["Enhanced", "Optimized z, OLS, 10bps", "2.8", "57%", "0.166"],
    ["Kalman", "Optimized z, Kalman \u03b2, 10bps", "2.6", "64%", "0.895"],
]
add_table(slide, 1.5, 1.6, 10, 2.0, headers, rows, col_widths=[1.5, 3.5, 1.5, 1.5, 1.5])

# Kalman fix explanation (left)
text_box(slide, 0.8, 3.9, 5.5, 0.5, "Kalman Terminal Beta Fix", font_size=20, bold=True, color=QUEENS_BLUE)

# Problem with adaptive
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5),
                                Inches(5.5), Inches(1.1))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.color.rgb = QUEENS_RED
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Adaptive Kalman (beta updated each OOS step)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = QUEENS_RED
p.font.name = 'Calibri'
p2 = tf.add_paragraph()
p2.text = "Prediction error \u2248 white noise by construction \u2192 artificial mean-reversion"
p2.font.size = Pt(12)
p2.font.color.rgb = DARK_GRAY
p2.font.name = 'Calibri'

# Fix
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8),
                                Inches(5.5), Inches(1.1))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_GREEN
shape.line.color.rgb = GREEN
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Fix: Kalman on calibration only \u2192 terminal beta"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = 'Calibri'
p2 = tf.add_paragraph()
p2.text = "Fixed beta applied to OOS = honest test of mean-reversion"
p2.font.size = Pt(12)
p2.font.color.rgb = DARK_GRAY
p2.font.name = 'Calibri'

# Top enhanced pairs (right)
text_box(slide, 7, 3.9, 5.5, 0.5, "Top Enhanced Pairs", font_size=20, bold=True, color=QUEENS_BLUE)
headers = ["Pair", "Type", "Enhanced PnL", "Kalman PnL"]
rows = [
    ["AMAT-TSM", "intra", "+30.30", "+54.19"],
    ["CVX-SU", "intra", "+13.16", "+10.41"],
    ["APLD-IREN", "cross", "+10.40", "+9.47"],
    ["RTX-SHEL", "cross", "+10.07", "+13.87"],
    ["HUT-RIOT", "intra", "+9.91", "+5.65"],
]
add_table(slide, 7, 4.5, 5.5, 2.8, headers, rows, col_widths=[1.5, 0.8, 1.5, 1.5])

add_slide_number(slide)
add_speaker_notes(slide, "The Kalman strategy wins: 64% profitable even with 10bps costs (vs 54% baseline without costs). Critical fix: an earlier version used adaptive Kalman beta during OOS, which creates artificial mean-reversion because the filter's prediction error is white noise by construction. The terminal-beta approach freezes the hedge ratio at the end of calibration and applies it as a constant to OOS data. Z-score optimization is a grid search over entry/exit thresholds and lookback on calibration data only.")


# ============================================================================
# SLIDE 15: Walk-Forward Validation
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Walk-Forward Validation",
         font_size=28, bold=True, color=QUEENS_BLUE)

text_box(slide, 0.8, 1.2, 11, 0.6,
         "5-split rolling calibration/OOS windows test whether performance persists across multiple temporal splits.",
         font_size=16, color=DARK_GRAY)

headers = ["Pair", "Splits", "Avg Sharpe", "Std Sharpe", "Total Trades"]
rows = [
    ["AFRM\u2013ZETA", "5", "2.07", "1.13", "11"],
    ["ENPH\u2013RUN", "5", "1.82", "0.82", "6"],
    ["COP\u2013CVE", "5", "0.45", "0.57", "16"],
]
add_table(slide, 2.5, 2.2, 8, 2.2, headers, rows, col_widths=[1.8, 1.0, 1.5, 1.5, 1.5])

# Two-column comparison
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.8),
                                Inches(5.0), Inches(2.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_GREEN
shape.line.color.rgb = GREEN
text_box(slide, 1.3, 4.9, 4.5, 0.4, "AFRM\u2013ZETA: Robust Signal", font_size=17, bold=True, color=GREEN)
bullet_slide(slide, 1.3, 5.4, 4.5, 1.5, [
    "Sharpe 2.07 with moderate variance (std 1.13)",
    "Consistent across all 5 temporal splits",
    "11 total trades \u2014 statistically meaningful",
], font_size=13, color=DARK_GRAY)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.8),
                                Inches(5.0), Inches(2.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_GREEN
shape.line.color.rgb = GREEN
text_box(slide, 7.1, 4.9, 4.5, 0.4, "ENPH\u2013RUN: Consistent", font_size=17, bold=True, color=GREEN)
bullet_slide(slide, 7.1, 5.4, 4.5, 1.5, [
    "Sharpe 1.82 with low variance (std 0.82)",
    "Solar energy pair \u2014 strong economic link",
    "6 trades \u2014 moderate confidence",
], font_size=13, color=DARK_GRAY)

add_slide_number(slide)
add_speaker_notes(slide, "Walk-forward validation addresses the single-split limitation of the main backtest. AFRM-ZETA shows robust performance with Sharpe 2.07 persisting across all 5 temporal splits. ENPH-RUN has low variance (std 0.82) reflecting consistent solar energy sector coupling. COP-CVE included as contrast with lower Sharpe (0.45) but high trade count (16).")


# ============================================================================
# SLIDE 16: Permutation Test Results
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Statistical Significance", "1,148 pairs confirmed by permutation testing")
add_speaker_notes(slide, "The permutation test validates that co-clustering structure is non-random.")


# ============================================================================
# SLIDE 17: Permutation Details
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Permutation Test: 1,148 Significant Pairs",
         font_size=28, bold=True, color=QUEENS_BLUE)

text_box(slide, 0.8, 1.2, 11, 0.8,
         "Feature-shuffle test: at 80 sampled timestamps \u00d7 30 permutations, shuffle feature vectors across tickers "
         "(preserving cross-feature correlation) and re-run the full pipeline. Compare observed co-clustering to null distribution. "
         "Primary threshold shown is Z > 1.645 (one-sided p < 0.05); stricter Z > 1.96 yields 940 pairs.",
         font_size=15, color=DARK_GRAY)

# Summary stats in boxes
stats = [
    ("3,539", "pairs tested", QUEENS_BLUE),
    ("1,148", "significant (Z > 1.645)", GREEN),
    ("671", "significant AND tradeable", QUEENS_RED),
]
for i, (val, label, color) in enumerate(stats):
    x = 1.0 + i * 4.0
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(2.3), Inches(3.5), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = color
    text_box(slide, x + 0.2, 2.35, 1.2, 0.5, val, font_size=24, bold=True, color=color)
    text_box(slide, x + 1.3, 2.4, 2.0, 0.5, label, font_size=13, color=DARK_GRAY)

# Top significant pairs
headers = ["Pair", "Z-Score", "Sector"]
rows = [
    ["PBR\u2013PBR-A", "38.59", "Energy"],
    ["AAL\u2013DAL", "38.19", "Industrials"],
    ["AAL\u2013UAL", "36.83", "Industrials"],
    ["DAL\u2013UAL", "34.03", "Industrials"],
    ["HAL\u2013SLB", "32.08", "Energy"],
    ["KMI\u2013WMB", "30.65", "Energy"],
    ["DVN\u2013OXY", "27.04", "Energy"],
    ["CVX\u2013XOM", "26.86", "Energy"],
    ["SU\u2013XOM", "25.86", "Energy"],
    ["AMAT\u2013LRCX", "25.13", "Technology"],
]
add_table(slide, 3, 3.3, 7, 3.8, headers, rows, col_widths=[2.5, 1.5, 3.0])

add_slide_number(slide)
add_speaker_notes(slide, "The permutation test shuffles which ticker has which features, breaking any real clustering structure. Under the null hypothesis, ~5% of 3,539 pairs (~177) would exceed Z > 1.645 by chance. The observed 1,148 is 6.5x this, indicating genuine signal. Top pairs (PBR-PBR.A, airlines, oil majors) have Z-scores far beyond any false discovery threshold. The Phase 1 pair AMAT-LRCX also appears (Z=25.13), confirming it generalizes to the larger universe.")


# ============================================================================
# SLIDE 18: What We Learned
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "What We Learned",
         font_size=28, bold=True, color=QUEENS_BLUE)

# What works
text_box(slide, 0.8, 1.3, 5.5, 0.5, "What Works", font_size=20, bold=True, color=GREEN)
bullet_slide(slide, 0.8, 1.9, 5.8, 3.5, [
    "Clustering discovers real economic relationships",
    "  \u2014 semiconductor subsectors, crypto mining clusters",
    "5-test framework: 59% tradeable vs 0% cointegration",
    "Kalman hedge: 64% profitable (vs 54% OLS baseline)",
    "Walk-forward validates top pairs across 5 splits",
    "1,148 pairs confirmed by permutation test",
], font_size=15, color=DARK_GRAY)

# Limitations
text_box(slide, 7, 1.3, 5.5, 0.5, "Limitations & Honest Findings", font_size=20, bold=True, color=QUEENS_RED)
bullet_slide(slide, 7, 1.9, 5.8, 3.5, [
    "Cross-sector pairs weaker (55.6% tradeable vs 64.1% intra)",
    "Short OOS (~68 days): 2\u20133 trades per pair",
    "Consensus \u2260 profitability (Phase 1 finding)",
    "10bps flat cost is simplified (no slippage/impact)",
    "No portfolio-level risk management",
    "Only 50 of 2,148 tradeable pairs met optimizer constraints (\u2265 5 calibration trades)",
], font_size=15, color=DARK_GRAY)

# Bottom line
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.5),
                                Inches(9), Inches(1.3))
shape.fill.solid()
shape.fill.fore_color.rgb = QUEENS_BLUE
shape.line.color.rgb = QUEENS_BLUE
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("The clustering methodology finds real structure \u2014 that's proven. The persistent pairs backtest\n"
          "validates the discovery engine. The real opportunity: trading transient formations as they happen.")
p.font.size = Pt(17)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide)
add_speaker_notes(slide, "Key framing: the current backtest of persistent pairs is a VALIDATION EXERCISE -- it proves the clustering finds real relationships. The unique edge of this methodology is in detecting transient formations (hours/days) that traditional traders don't monitor. The persistent pairs (airlines, energy) are a byproduct of similar fundamentals; the transient formations are where the unexploited signal lives. The statistical foundation (permutation testing, 5-test framework, Kalman hedge) is established and ready for the next step: event-driven transient trading.")


# ============================================================================
# SLIDE 19: Future Work
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
text_box(slide, 0.8, 0.4, 10, 0.8, "Future Work",
         font_size=28, bold=True, color=QUEENS_BLUE)

items = [
    ("\u2705", "Walk-Forward Validation", "5-split rolling cal/OOS", True),
    ("\u2705", "Adaptive Hedge Ratios", "Kalman terminal beta estimation", True),
    ("\u2705", "Expanded Universe", "142 tickers, 5 sectors via 3-layer screening", True),
    ("\u2705", "Transaction Costs", "10bps round-trip incorporated", True),
    ("\u2B50", "Transient Event Trading", "Trade formation/dissolution events directly \u2014 the core application", False),
    ("\u2B1C", "Real-Time Pipeline", "Live clustering \u2192 formation detection \u2192 signal generation", False),
    ("\u2B1C", "Portfolio Optimization", "Correlated pairs risk management", False),
    ("\u2B1C", "Longer Historical Data", "Multi-year hourly data for robust estimation", False),
]

for i, (status, title, desc, done) in enumerate(items):
    y = 1.3 + i * 0.7
    color = GREEN if done else DARK_GRAY
    text_box(slide, 1.0, y, 0.5, 0.5, status, font_size=18, color=color)
    text_box(slide, 1.6, y, 4, 0.5, title, font_size=16, bold=True, color=QUEENS_BLUE)
    text_box(slide, 5.8, y, 6.5, 0.5, desc, font_size=14, color=DARK_GRAY)

add_slide_number(slide)
add_speaker_notes(slide, "Phase 2 completed the first four items from Phase 1's future work list. The star item -- transient event trading -- is the natural next step: using cluster formation/dissolution events as trade signals instead of static z-score strategies on persistent pairs. The formation/dissolution infrastructure already exists in the codebase. The remaining items support production deployment.")


# ============================================================================
# SLIDE 20: Thank You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, QUEENS_BLUE)

text_box(slide, 2, 2.5, 9, 1.5,
         "Thank You",
         font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

text_box(slide, 2, 4.2, 9, 0.8,
         "Questions?",
         font_size=32, color=GOLD, alignment=PP_ALIGN.CENTER)

text_box(slide, 2, 5.5, 9, 0.6,
         "Queen's Quant Club  |  February 2026",
         font_size=18, color=HEADER_BLUE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide)
add_speaker_notes(slide, "Thank the audience and open the floor for questions.")


# ============================================================================
# SAVE
# ============================================================================
import os
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentation.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
